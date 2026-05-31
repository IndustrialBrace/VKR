#!/usr/bin/env python3
"""Конвертация PPTX -> Markdown с извлечением текста, таблиц и заметок."""
import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def shape_to_md(shape, lines):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            shape_to_md(s, lines)
        return

    if shape.has_table:
        table = shape.table
        rows = list(table.rows)
        if not rows:
            return
        header = [cell.text.strip().replace("\n", " ") for cell in rows[0].cells]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
        return

    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                # fallback на para.text если runs пустые
                text = para.text.strip()
            if not text:
                continue
            level = para.level if para.level is not None else 0
            indent = "  " * level
            lines.append(f"{indent}- {text}")
        lines.append("")


def convert(pptx_path, md_path):
    prs = Presentation(pptx_path)
    out = []
    out.append(f"# {os.path.splitext(os.path.basename(pptx_path))[0]}")
    out.append("")
    out.append(f"_Преобразовано из {os.path.basename(pptx_path)}_")
    out.append("")

    for i, slide in enumerate(prs.slides, 1):
        # Заголовок слайда
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        out.append(f"## Слайд {i}" + (f". {title}" if title else ""))
        out.append("")

        title_id = id(slide.shapes.title) if slide.shapes.title else None
        for shape in slide.shapes:
            if id(shape) == title_id:
                continue
            shape_to_md(shape, out)

        # Заметки докладчика
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = (notes_tf.text or "").strip() if notes_tf else ""
            if notes_text:
                out.append("**Заметки к слайду:**")
                out.append("")
                for line in notes_text.splitlines():
                    line = line.strip()
                    if line:
                        out.append(f"> {line}")
                out.append("")

        out.append("---")
        out.append("")

    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(out))


if __name__ == "__main__":
    src = sys.argv[1]
    dst = sys.argv[2]
    convert(src, dst)
    print(f"OK: {dst}")
